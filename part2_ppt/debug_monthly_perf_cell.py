from __future__ import annotations

import sys
from pathlib import Path
from typing import Optional, Tuple

from pptx import Presentation


def norm(s: str) -> str:
    return (s or "").replace("\r", "").replace(" \n", "\n").strip()


def find_table_shape(prs: Presentation, shape_name: str):
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "name", None) == shape_name:
                return slide_idx, shape
    return None, None


def find_col_by_header(tbl, header_text: str) -> Optional[int]:
    for c in range(len(tbl.columns)):
        if norm(tbl.cell(1, c).text) == header_text:
            return c
    return None


def find_row_by_token(tbl, col_idx: int, token: str) -> Optional[int]:
    for r in range(2, len(tbl.rows)):
        txt = tbl.cell(r, col_idx).text or ""
        if token in txt:
            return r
    return None


def safe_run_font_snapshot(cell) -> str:
    try:
        tf = cell.text_frame
        if not tf.paragraphs:
            return "no_paragraphs"
        p0 = tf.paragraphs[0]
        align = str(p0.alignment)
        if not p0.runs:
            return f"alignment={align} runs=0"
        r0 = p0.runs[0]
        f = r0.font
        rgb = None
        try:
            if f.color is not None and f.color.rgb is not None:
                rgb = str(f.color.rgb)
        except Exception:
            rgb = "error"
        return (
            f"alignment={align} runs={len(p0.runs)} "
            f"name={f.name} size={f.size} bold={f.bold} italic={f.italic} underline={f.underline} color={rgb}"
        )
    except Exception as e:
        return f"error_snapshot: {e}"


def extract_txbody_xml(cell) -> str:
    try:
        txBody = cell._tc.txBody
        return txBody.xml
    except Exception as e:
        return f"<error>{e}</error>"


def find_xml_markers(xml: str) -> str:
    markers = []
    if "<a:pPr" in xml:
        markers.append("has_a:pPr")
    else:
        markers.append("no_a:pPr")

    if "<a:endParaRPr" in xml:
        markers.append("has_a:endParaRPr")
    else:
        markers.append("no_a:endParaRPr")

    if "<a:r>" in xml or "<a:r " in xml:
        markers.append("has_a:r")
    else:
        markers.append("no_a:r")

    if "<a:latin" in xml:
        markers.append("has_a:latin")
    else:
        markers.append("no_a:latin")

    if "srgbClr" in xml:
        markers.append("has_srgbClr")
    else:
        markers.append("no_srgbClr")

    return ", ".join(markers)


def write_text(path: Path, text: str) -> None:
    path.write_text(text, encoding="utf-8")


def compare_cells(
    prs_path: Path,
    shape_name: str,
    month_year_header: str,
    target_token: str,
    control_token: str,
) -> None:
    prs = Presentation(str(prs_path))

    slide_idx, shape = find_table_shape(prs, shape_name)
    if shape is None:
        print(f"Could not find shape '{shape_name}' in {prs_path}")
        return

    if not hasattr(shape, "table"):
        print(f"Shape '{shape_name}' found on slide {slide_idx} but it is not a table.")
        return

    tbl = shape.table

    col_month_year = find_col_by_header(tbl, month_year_header)
    if col_month_year is None:
        print(f"Could not find Month Year column using header '{month_year_header}' on row 2.")
        return

    target_row = 11   # bad formatting row (PowerPoint row 12)
    control_row = 10  # good formatting row

    if target_row is None:
        print(f"Could not find a row whose Month Year cell contains '{target_token}'.")
        return
    if control_row is None:
        print(f"Could not find a row whose Month Year cell contains '{control_token}'.")
        return

    target_cell = tbl.cell(target_row, col_month_year)
    control_cell = tbl.cell(control_row, col_month_year)

    print(f"PPTX: {prs_path}")
    print(f"Shape: {shape_name} on slide {slide_idx}")
    print(f"Month Year column index: {col_month_year}")
    print(f"Target row for {target_token}: {target_row}")
    print(f"Control row for {control_token}: {control_row}")
    print()

    print("Text values")
    print(f"Target cell text: {repr(target_cell.text)}")
    print(f"Control cell text: {repr(control_cell.text)}")
    print()

    print("python-pptx visible formatting snapshot")
    print(f"Target:  {safe_run_font_snapshot(target_cell)}")
    print(f"Control: {safe_run_font_snapshot(control_cell)}")
    print()

    target_xml = extract_txbody_xml(target_cell)
    control_xml = extract_txbody_xml(control_cell)

    print("XML markers")
    print(f"Target:  {find_xml_markers(target_xml)}")
    print(f"Control: {find_xml_markers(control_xml)}")
    print()

    out_dir = prs_path.parent
    target_xml_path = out_dir / f"debug_{shape_name}_{target_token.strip('[]')}_txBody.xml"
    control_xml_path = out_dir / f"debug_{shape_name}_{control_token.strip('[]')}_txBody.xml"

    write_text(target_xml_path, target_xml)
    write_text(control_xml_path, control_xml)

    print("Wrote XML files")
    print(str(target_xml_path))
    print(str(control_xml_path))
    print()
    print("If Target lacks a:pPr or a:endParaRPr while Control has them, that is the root cause.")
    print("If both have them, the next thing to check is run structure and font nodes in the XML.")


def main() -> None:
    if len(sys.argv) < 2:
        print("Usage: python debug_monthly_perf_cell.py <path_to_pptx>")
        print("Example: python debug_monthly_perf_cell.py H:\\path\\to\\__tmp_updated_...pptx")
        return

    prs_path = Path(sys.argv[1])
    if not prs_path.exists():
        print(f"File not found: {prs_path}")
        return

    compare_cells(
        prs_path=prs_path,
        shape_name="monthly_perf_table",
        month_year_header="Month Year",
        target_token="[T4]",
        control_token="[T5]",
    )


if __name__ == "__main__":
    main()
