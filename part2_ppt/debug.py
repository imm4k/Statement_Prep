# debug.py
from __future__ import annotations

import os
import sys
from pathlib import Path
from typing import Optional, List, Tuple

from pptx import Presentation

import config


def _print_banner(title: str) -> None:
    print("\n" + "=" * 80)
    print(title)
    print("=" * 80)


def _norm(s: str) -> str:
    return (s or "").replace("\r", "").replace(" \n", "\n").strip()


def _guess_ppt_path() -> Optional[Path]:
    """
    Attempts to locate a PPTX path from config by scanning for common attribute names.
    Falls back to searching the working directory recursively for .pptx files.
    """
    candidates: List[str] = []
    for name in dir(config):
        if not name.isupper():
            continue
        if "PPT" in name or "PPTX" in name or "PRESENTATION" in name:
            candidates.append(name)

    preferred_names = [
        "OUTPUT_PPTX_PATH",
        "OUTPUT_PPT_PATH",
        "PPTX_OUTPUT_PATH",
        "PPT_OUTPUT_PATH",
        "TEMPLATE_PPTX_PATH",
        "TEMPLATE_PPT_PATH",
        "PPTX_TEMPLATE_PATH",
        "PPT_TEMPLATE_PATH",
        "PRESENTATION_PATH",
        "PPTX_PATH",
    ]

    ordered: List[str] = []
    for n in preferred_names:
        if hasattr(config, n):
            ordered.append(n)
    for n in candidates:
        if n not in ordered:
            ordered.append(n)

    for attr in ordered:
        try:
            v = getattr(config, attr)
        except Exception:
            continue
        if v is None:
            continue
        try:
            p = Path(str(v)).expanduser()
        except Exception:
            continue
        if p.exists() and p.suffix.lower() == ".pptx":
            print(f"Using config.{attr} -> {p}")
            return p

    # Fallback: search current folder for .pptx
    here = Path(__file__).resolve().parent
    pptxs = sorted(here.rglob("*.pptx"))
    if pptxs:
        print("No PPTX path found in config. Found PPTX files under this folder:")
        for i, p in enumerate(pptxs[:25], start=1):
            print(f"  {i:>2}. {p}")
        if len(pptxs) > 25:
            print(f"  ... plus {len(pptxs) - 25} more")
        return None

    return None


def _print_table_preview(tbl, max_rows: int = 6, max_cols: int = 10) -> None:
    rows = len(tbl.rows)
    cols = len(tbl.columns)
    print(f"Table size: rows={rows}, cols={cols}")

    r_lim = min(rows, max_rows)
    c_lim = min(cols, max_cols)

    for r in range(r_lim):
        cells = []
        for c in range(c_lim):
            t = _norm(tbl.cell(r, c).text)
            t = t.replace("\n", "\\n")
            if len(t) > 40:
                t = t[:37] + "..."
            cells.append(t)
        print(f"row {r}: {cells}")

    if rows > r_lim:
        print(f"... {rows - r_lim} more rows not shown")
    if cols > c_lim:
        print(f"... {cols - c_lim} more cols not shown")


def _find_shapes(prs: Presentation, target_name: str) -> List[Tuple[int, object]]:
    hits: List[Tuple[int, object]] = []
    for si, slide in enumerate(prs.slides):
        for shp in slide.shapes:
            name = getattr(shp, "name", "")
            if str(name) == target_name:
                hits.append((si, shp))
    return hits


def _find_near_matches(prs: Presentation, substrings: List[str]) -> List[Tuple[int, str, bool]]:
    out: List[Tuple[int, str, bool]] = []
    for si, slide in enumerate(prs.slides):
        for shp in slide.shapes:
            nm = str(getattr(shp, "name", ""))
            nm_low = nm.lower()
            if all(s.lower() in nm_low for s in substrings):
                out.append((si, nm, hasattr(shp, "table")))
    return out


def main() -> None:
    _print_banner("debug.py: available_cash diagnostics")

    print(f"Python: {sys.version.split()[0]}")
    print(f"CWD: {os.getcwd()}")
    print(f"Script dir: {Path(__file__).resolve().parent}")

    _print_banner("Config scan for PPTX paths")
    ppt_path = _guess_ppt_path()
    if ppt_path is None:
        print("\nCould not auto-detect which PPTX file to open.")
        print("Paste a full path to the PPTX you expect to be updating, then press Enter.")
        user_in = input("PPTX path: ").strip().strip('"').strip("'")
        if not user_in:
            print("No path provided. Exiting.")
            return
        ppt_path = Path(user_in).expanduser()
        if not ppt_path.exists():
            print(f"Path not found: {ppt_path}")
            return
        if ppt_path.suffix.lower() != ".pptx":
            print(f"Not a .pptx file: {ppt_path}")
            return

    _print_banner("Open PPTX")
    print(f"Opening: {ppt_path}")
    prs = Presentation(str(ppt_path))
    print(f"Slides: {len(prs.slides)}")

    target = "available_cash"

    _print_banner(f"Search for shape name exactly: {target!r}")
    hits = _find_shapes(prs, target)
    if not hits:
        print("NOT FOUND: No shape has name exactly 'available_cash'.")
        print("\nNear-match search (contains 'available' and 'cash'):")
        near = _find_near_matches(prs, ["available", "cash"])
        if not near:
            print("No near matches found either.")
        else:
            for si, nm, has_tbl in near:
                print(f"  slide {si+1}: name={nm!r}, has_table={has_tbl}")
        print("\nMost likely cause: the table shape name in PowerPoint is not 'available_cash'.")
        return

    print(f"FOUND: {len(hits)} shape(s) named 'available_cash'")
    for si, shp in hits:
        print(f"\nslide {si+1}:")
        print(f"  shape id: {getattr(shp, 'shape_id', None)}")
        print(f"  name: {getattr(shp, 'name', '')!r}")
        print(f"  has_table: {hasattr(shp, 'table')}")
        print(f"  shape_type: {getattr(shp, 'shape_type', None)}")

        if not hasattr(shp, "table"):
            print("  This shape is not a table. Updater will exit early.")
            continue

        tbl = shp.table
        _print_table_preview(tbl, max_rows=6, max_cols=6)

        _print_banner("Header row detection preview")
        for candidate in (0, 1):
            if candidate >= len(tbl.rows):
                continue
            headers = [_norm(tbl.cell(candidate, c).text) for c in range(len(tbl.columns))]
            print(f"Row {candidate} headers:")
            for idx, h in enumerate(headers):
                print(f"  col {idx}: {h!r}")

    _print_banner("Done")
    print("If the shape exists and is a table but still does not update, the issue is likely:")
    print("1) OBJECT_UPDATERS does not include 'available_cash', or")
    print("2) The update engine is not dispatching by shape.name the way you expect, or")
    print("3) You are running against a different PPTX than the one being inspected here.")


if __name__ == "__main__":
    main()
