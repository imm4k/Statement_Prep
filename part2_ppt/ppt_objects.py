from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime
from typing import Callable, Dict, Optional

from pptx.presentation import Presentation
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

@dataclass(frozen=True)
class UpdateContext:
    investor: str
    owner: Optional[str]
    ownership_pct: float
    ownership_factor: float
    statement_thru_date_dt: datetime
    statement_thru_date_str: str
    t1_str: str

ObjectUpdater = Callable[[Slide, BaseShape, Presentation, UpdateContext], None]

def apply_ownership_amount(ctx: UpdateContext, amount: float, key: str) -> float:
    import config

    if bool(getattr(config, "OWNERSHIP_FORCE_100_PCT_IN_PART2", False)):
        return float(amount or 0.0)

    if ctx.ownership_pct >= 100.0:
        return float(amount or 0.0)

    k = str(key or "").strip()
    if k != "" and k in getattr(config, "OWNERSHIP_SCALING_EXCEPTIONS", []):
        return float(amount or 0.0)

    return float(amount or 0.0) * float(ctx.ownership_factor or 0.0)

def apply_object_updates(prs: Presentation, ctx: UpdateContext) -> None:
    from ppt_object_logic import OBJECT_UPDATERS

    shape_names = list(OBJECT_UPDATERS.keys())
    if not shape_names:
        print("No object updaters registered yet. Skipping object updates.")
        return

    total_hits = 0
    slides_count = len(prs.slides)
    print(f"Starting object updates across {slides_count} slides. Objects: {len(shape_names)}")

    for slide_idx, slide in enumerate(prs.slides, start=1):
        if slide_idx % 5 == 0 or slide_idx == 1 or slide_idx == slides_count:
            print(f"Scanning slide {slide_idx} of {slides_count}")

        for shape in slide.shapes:
            if not shape.name:
                continue
            updater = OBJECT_UPDATERS.get(shape.name)
            if updater is None:
                continue

            total_hits += 1
            print(f"Updating object '{shape.name}' on slide {slide_idx} for investor '{ctx.investor}'")
            updater(slide, shape, prs, ctx)

    print(f"Completed object updates. Objects updated: {total_hits}")
