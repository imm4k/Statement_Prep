from __future__ import annotations

import sqlite3
from datetime import datetime
from typing import Dict, Tuple

import config

from pptx.presentation import Presentation
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from ppt_objects import UpdateContext
from ppt_objects import apply_ownership_amount


def _owner_filter_sql(ctx: UpdateContext) -> tuple[str, tuple]:
    if ctx.owner is None or str(ctx.owner).strip() == "":
        return "", tuple()
    return " AND owner = ? ", (str(ctx.owner).strip(),)


def _coerce_date_yyyy_mm_dd(value: object) -> datetime:
    if isinstance(value, datetime):
        return value
    s = str(value or "").strip()
    if s == "":
        raise ValueError("Empty acquired value")
    s = s.replace("T", " ")
    s10 = s[:10]
    return datetime.strptime(s10, "%Y-%m-%d")

def _set_cell_text_preserve_cell_format(cell, text: str) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    txBody = cell._tc.txBody
    p = txBody.find("a:p", ns)
    if p is None:
        cell.text_frame.text = text
        return

    pPr = p.find("a:pPr", ns)
    endParaRPr = p.find("a:endParaRPr", ns)

    algn = None
    if pPr is not None:
        algn = pPr.get("algn")

    typeface = None
    sz = None
    b = None
    i = None
    u = None
    color_val = None

    if endParaRPr is not None:
        sz = endParaRPr.get("sz")
        b = endParaRPr.get("b")
        i = endParaRPr.get("i")
        u = endParaRPr.get("u")

        latin = endParaRPr.find("a:latin", ns)
        if latin is not None:
            typeface = latin.get("typeface")

        srgb = endParaRPr.find(".//a:srgbClr", ns)
        if srgb is not None:
            color_val = srgb.get("val")

    cell.text_frame.text = text

    p0 = cell.text_frame.paragraphs[0]
    if algn == "ctr":
        p0.alignment = PP_ALIGN.CENTER
    elif algn == "l":
        p0.alignment = PP_ALIGN.LEFT
    elif algn == "r":
        p0.alignment = PP_ALIGN.RIGHT
    elif algn == "just":
        p0.alignment = PP_ALIGN.JUSTIFY

    if p0.runs:
        r0 = p0.runs[0]
        if typeface:
            r0.font.name = typeface
        if sz and str(sz).isdigit():
            r0.font.size = Pt(int(sz) / 100)
        if b is not None:
            r0.font.bold = (str(b) == "1")
        if i is not None:
            r0.font.italic = (str(i) == "1")
        if u is not None:
            r0.font.underline = (str(u).lower() != "none")
        if color_val and len(color_val) == 6:
            r0.font.color.rgb = RGBColor(
                int(color_val[0:2], 16),
                int(color_val[2:4], 16),
                int(color_val[4:6], 16),
            )


def update_summary_table(slide: Slide, shape: BaseShape, prs: Presentation, ctx: UpdateContext) -> None:
    if not hasattr(shape, "table"):
        return

    from openpyxl import load_workbook
    from pptx.dml.color import RGBColor

    def _norm_header(s: str) -> str:
        return s.replace("\r", "").replace(" \n", "\n").strip()

    def _fmt_currency(x: float) -> tuple[str, bool]:
        if abs(x) < 0.5:
            return "-", False
        if x < 0:
            return f"(${abs(x):,.0f})", True
        return f"${x:,.0f}", False

    def _apply_red_if_negative(cell, is_negative: bool) -> None:
        if not is_negative:
            return
        try:
            p0 = cell.text_frame.paragraphs[0]
            if p0.runs:
                p0.runs[0].font.color.rgb = RGBColor(255, 0, 0)
        except Exception:
            return

    def _set_currency_cell(cell, amount: float) -> None:
        txt, is_neg = _fmt_currency(amount)
        _set_cell_text_preserve_cell_format(cell, txt)
        _apply_red_if_negative(cell, is_neg)

    def _read_general_config_market_values() -> Dict[str, float]:
        labels = {
            "Studio": "Studio Market:",
            "1-Bed": "1-Bed Market:",
            "2-Bed": "2-Bed Market:",
            "3-Bed": "3-Bed Market:",
        }

        wb = load_workbook(filename=str(config.SETUP_EXCEL_PATH), data_only=True)
        try:
            if config.GENERAL_CONFIG_SHEET not in wb.sheetnames:
                raise RuntimeError(f"Missing sheet: {config.GENERAL_CONFIG_SHEET}")

            ws = wb[config.GENERAL_CONFIG_SHEET]
            label_to_value: Dict[str, float] = {}

            for r in range(1, ws.max_row + 1):
                a = ws.cell(r, 1).value
                b = ws.cell(r, 2).value
                if a is None:
                    continue
                key = str(a).strip()
                if key in labels.values():
                    try:
                        label_to_value[key] = float(b)
                    except Exception:
                        label_to_value[key] = 0.0

            out: Dict[str, float] = {}
            for unit_type, label in labels.items():
                out[unit_type] = float(label_to_value.get(label, 0.0))

            return out
        finally:
            wb.close()

    tbl = shape.table

    col_property = None
    col_type = None
    col_est_mkt_value = None
    col_mortgage_balance = None
    col_nav = None

    for c in range(len(tbl.columns)):
        header = _norm_header(tbl.cell(1, c).text)

        if header == "Property":
            col_property = c
        elif header == "Type":
            col_type = c
        elif header == "Estimated\nMarket Value":
            col_est_mkt_value = c
        elif header == "Mortgage\nBalance":
            col_mortgage_balance = c
        elif header in ("Net Asset Value (NAV)", "Net Asset\nValue (NAV)"):
            col_nav = c

    if col_property is None:
        print("summary_table missing required column header: Property")
        return

    total_row_idx = len(tbl.rows) - 1

    market_values_by_type = _read_general_config_market_values()

    con = sqlite3.connect(str(config.SQLITE_PATH))

    owner_sql, owner_params = _owner_filter_sql(ctx)

    sql_mortgage_balance = f"""
        SELECT property,
               ABS(SUM(value)) AS mortgage_balance
        FROM gl_agg
        WHERE investor = ?
          AND categorization = 'Mortgage Balance'
          AND (timeframe IS NULL OR timeframe <> 'N/A')
          {owner_sql}
        GROUP BY property
    """

    mortgage_rows = con.execute(sql_mortgage_balance, (ctx.investor, *owner_params)).fetchall()

    con.close()

    mortgage_by_prop: Dict[str, float] = {}
    for prop, v in mortgage_rows:
        if prop is None:
            continue
        mortgage_by_prop[str(prop)] = float(v or 0.0)

    est_hits = 0
    mortgage_hits = 0
    nav_hits = 0

    est_total = 0.0
    mortgage_total = 0.0
    nav_total = 0.0

    data_row_count = max(0, len(tbl.rows) - 3)
    print(f"summary_table Starting process for {data_row_count} rows.")

    current = 0
    for r in range(2, len(tbl.rows)):
        if r == total_row_idx:
            continue

        current += 1
        print(f"summary_table Currently on {current} of {data_row_count}")

        prop_name = tbl.cell(r, col_property).text.strip()
        if prop_name == "":
            continue

        unit_type = ""
        if col_type is not None:
            unit_type = tbl.cell(r, col_type).text.strip()

        if unit_type == "" and col_type is not None:
            continue

        est_mkt_raw = float(market_values_by_type.get(unit_type, 0.0))
        mortgage_bal_raw = abs(mortgage_by_prop.get(prop_name, 0.0))
        nav_raw = est_mkt_raw - mortgage_bal_raw

        est_mkt = apply_ownership_amount(ctx, est_mkt_raw, "summary_table.estimated_market_value")
        mortgage_bal = apply_ownership_amount(ctx, mortgage_bal_raw, "summary_table.mortgage_balance")
        nav = apply_ownership_amount(ctx, nav_raw, "summary_table.nav")

        if col_est_mkt_value is not None:
            _set_currency_cell(tbl.cell(r, col_est_mkt_value), est_mkt)
            est_hits += 1

        if col_mortgage_balance is not None:
            _set_currency_cell(tbl.cell(r, col_mortgage_balance), mortgage_bal)
            mortgage_hits += 1

        if col_nav is not None:
            _set_currency_cell(tbl.cell(r, col_nav), nav)
            nav_hits += 1

        est_total += est_mkt
        mortgage_total += mortgage_bal
        nav_total += nav

    if col_est_mkt_value is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_est_mkt_value), est_total)

    if col_mortgage_balance is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_mortgage_balance), mortgage_total)

    if col_nav is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_nav), nav_total)

    print(f"summary_table estimated_market_value updated rows: {est_hits}")
    print(f"summary_table mortgage_balance updated rows: {mortgage_hits}")
    print(f"summary_table nav updated rows: {nav_hits}")

def update_monthly_perf_table(slide: Slide, shape: BaseShape, prs: Presentation, ctx: UpdateContext) -> None:
    if not hasattr(shape, "table"):
        return

    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    from ppt_monthly_stmt_values import build_month_year_labels

    owner_sql, owner_params = _owner_filter_sql(ctx)

    def _norm_header(s: str) -> str:
        return (s or "").replace("\r", "").replace(" \n", "\n").strip()

    def _fmt_currency(x: float) -> tuple[str, bool]:
        if abs(x) < 0.5:
            return "-", False
        if x < 0:
            return f"(${abs(x):,.0f})", True
        return f"${x:,.0f}", False

    def _apply_red_if_negative(cell, is_negative: bool) -> None:
        if not is_negative:
            return
        try:
            p0 = cell.text_frame.paragraphs[0]
            if p0.runs:
                p0.runs[0].font.color.rgb = RGBColor(255, 0, 0)
        except Exception:
            return

    def _set_currency_cell(cell, amount: float) -> None:
        txt, is_neg = _fmt_currency(amount)
        _set_cell_text_preserve_cell_format(cell, txt)
        _apply_red_if_negative(cell, is_neg)

    tbl = shape.table

    col_month_year = None
    col_rent = None
    col_dividend = None
    col_total_rev = None
    col_hoa_mgt = None
    col_repairs_other = None
    col_mortgage_int = None
    col_total_exp = None
    col_monthly = None
    col_cumulative = None

    def _hdr_compact(s: str) -> str:
        return _norm_header(s).replace("\n", " ").strip()

    for c in range(len(tbl.columns)):
        header = _norm_header(tbl.cell(1, c).text)
        header_compact = _hdr_compact(tbl.cell(1, c).text)

        if header == "Month Year":
            col_month_year = c
        elif header == "Rent":
            col_rent = c
        elif header == "Dividend":
            col_dividend = c
        elif header == "Total Revenue":
            col_total_rev = c
        elif header_compact == "HOA & Mgt. Fee":
            col_hoa_mgt = c
        elif header_compact == "Repairs & Other Exp.":
            col_repairs_other = c
        elif header_compact == "Mortgage Interest":
            col_mortgage_int = c
        elif header == "Total Expenses":
            col_total_exp = c
        elif header == "Monthly":
            col_monthly = c
        elif header == "Cumulative":
            col_cumulative = c

    if col_month_year is None:
        print("monthly_perf_table missing required column header: Month Year")
        return

    token_to_label = build_month_year_labels(ctx, property_name=None)

    wanted_cats = (
        "Rent",
        "Dividend",
        "HOA & Mgt. Fee",
        "Repairs & Other Exp.",
        "Mortgage Interest",
    )
    placeholders = ",".join(["?"] * len(wanted_cats))

    sql = f"""
            SELECT timeframe, categorization, gl_mapping_type, SUM(value) AS total_value
            FROM gl_agg
            WHERE investor = ?
            AND (timeframe IS NULL OR timeframe <> 'N/A')
            AND timeframe IN ('[T1]','[T2]','[T3]','[T4]','[T5]','[T6]','[T7]','[T8]','[T9]','[T10]','[T11]','[T12]','[T13]')
            AND categorization IN ({placeholders})
            {owner_sql}
            GROUP BY timeframe, categorization, gl_mapping_type
        """

    con = sqlite3.connect(str(config.SQLITE_PATH))
    rows = con.execute(sql, (ctx.investor, *wanted_cats, *owner_params)).fetchall()
    con.close()

    vals: Dict[str, Dict[str, float]] = {}
    for tf, cat, mapping_type, total_value in rows:
        tf_key = str(tf).strip()
        cat_key = str(cat).strip()
        mt = str(mapping_type or "").strip().lower()
        v = float(total_value or 0.0)

        if mt in ("revenue", "expense"):
            v = -1.0 * v

        if tf_key not in vals:
            vals[tf_key] = {}
        vals[tf_key][cat_key] = vals[tf_key].get(cat_key, 0.0) + v

    total_row_idx = len(tbl.rows) - 1
    cumulative_running = 0.0

    total_rent = 0.0
    total_dividend = 0.0
    total_hoa_mgt = 0.0
    total_repairs_other = 0.0
    total_mortgage_int = 0.0

    data_row_count = max(0, len(tbl.rows) - 3)
    print(f"monthly_perf_table Starting process for {data_row_count} rows.")

    current = 0
    for r in range(2, len(tbl.rows)):
        if r == total_row_idx:
            continue

        row_label = tbl.cell(r, col_month_year).text.strip()
        if row_label == "":
            continue

        tf_token = None
        for token in (f"[T{n}]" for n in range(1, 14)):
            if token in row_label:
                tf_token = token
                break

        if tf_token is None:
            continue

        if tf_token not in token_to_label:
            continue

        current += 1
        print(f"monthly_perf_table Currently on {current} of {data_row_count}")

        new_label = row_label.replace(tf_token, token_to_label[tf_token])
        _set_cell_text_preserve_cell_format(tbl.cell(r, col_month_year), new_label)
        p = tbl.cell(r, col_month_year).text_frame.paragraphs[0]
        if p.runs:
            r0 = p.runs[0]
            r0.font.name = "Lato"
            r0.font.size = Pt(10)
            r0.font.color.rgb = RGBColor(0, 0, 0)

        tf_vals = vals.get(tf_token, {})

        rent = apply_ownership_amount(ctx, float(tf_vals.get("Rent", 0.0)), "monthly_perf_table.rent")
        dividend = apply_ownership_amount(ctx, float(tf_vals.get("Dividend", 0.0)), "monthly_perf_table.dividend")

        hoa_mgt = apply_ownership_amount(ctx, float(tf_vals.get("HOA & Mgt. Fee", 0.0)), "monthly_perf_table.hoa_mgt_fee")
        repairs_other = apply_ownership_amount(ctx, float(tf_vals.get("Repairs & Other Exp.", 0.0)), "monthly_perf_table.repairs_other")
        mortgage_int = apply_ownership_amount(ctx, float(tf_vals.get("Mortgage Interest", 0.0)), "monthly_perf_table.mortgage_interest")

        total_rev = rent + dividend
        total_exp = hoa_mgt + repairs_other + mortgage_int
        monthly = total_rev + total_exp

        if cumulative_running == 0.0:
            cumulative_running = monthly
        else:
            cumulative_running = cumulative_running + monthly

        if col_rent is not None:
            _set_currency_cell(tbl.cell(r, col_rent), rent)
        if col_dividend is not None:
            _set_currency_cell(tbl.cell(r, col_dividend), dividend)
        if col_total_rev is not None:
            _set_currency_cell(tbl.cell(r, col_total_rev), total_rev)

        if col_hoa_mgt is not None:
            _set_currency_cell(tbl.cell(r, col_hoa_mgt), hoa_mgt)
        if col_repairs_other is not None:
            _set_currency_cell(tbl.cell(r, col_repairs_other), repairs_other)
        if col_mortgage_int is not None:
            _set_currency_cell(tbl.cell(r, col_mortgage_int), mortgage_int)
        if col_total_exp is not None:
            _set_currency_cell(tbl.cell(r, col_total_exp), total_exp)

        if col_monthly is not None:
            _set_currency_cell(tbl.cell(r, col_monthly), monthly)
        if col_cumulative is not None:
            _set_currency_cell(tbl.cell(r, col_cumulative), cumulative_running)

        total_rent += rent
        total_dividend += dividend
        total_hoa_mgt += hoa_mgt
        total_repairs_other += repairs_other
        total_mortgage_int += mortgage_int

    total_rev_all = total_rent + total_dividend
    total_exp_all = total_hoa_mgt + total_repairs_other + total_mortgage_int
    total_monthly_all = total_rev_all + total_exp_all

    if col_rent is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_rent), total_rent)
    if col_dividend is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_dividend), total_dividend)
    if col_total_rev is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_total_rev), total_rev_all)

    if col_hoa_mgt is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_hoa_mgt), total_hoa_mgt)
    if col_repairs_other is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_repairs_other), total_repairs_other)
    if col_mortgage_int is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_mortgage_int), total_mortgage_int)
    if col_total_exp is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_total_exp), total_exp_all)

    if col_monthly is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_monthly), total_monthly_all)
    if col_cumulative is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_cumulative), cumulative_running)

    print("monthly_perf_table updated.")

def update_monthly_cash_table(slide: Slide, shape: BaseShape, prs: Presentation, ctx: UpdateContext) -> None:
    if not hasattr(shape, "table"):
        return

    owner_sql, owner_params = _owner_filter_sql(ctx)

    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    from datetime import date

    from ppt_monthly_stmt_values import build_month_year_labels

    def _norm_header(s: str) -> str:
        return (s or "").replace("\r", "").replace(" \n", "\n").strip()

    def _hdr_compact(s: str) -> str:
        return _norm_header(s).replace("\n", " ").strip()

    def _fmt_currency(x: float) -> tuple[str, bool]:
        if abs(x) < 0.5:
            return "-", False
        if x < 0:
            return f"(${abs(x):,.0f})", True
        return f"${x:,.0f}", False

    def _apply_red_if_negative(cell, is_negative: bool) -> None:
        if not is_negative:
            return
        try:
            p0 = cell.text_frame.paragraphs[0]
            if p0.runs:
                p0.runs[0].font.color.rgb = RGBColor(255, 0, 0)
        except Exception:
            return

    def _set_currency_cell(cell, amount: float) -> None:
        txt, is_neg = _fmt_currency(amount)
        _set_cell_text_preserve_cell_format(cell, txt)
        _apply_red_if_negative(cell, is_neg)

    tbl = shape.table

    col_month_year = None
    col_owner_contrib = None
    col_mortgage_loan = None
    col_rent_dividend = None
    col_total_inflow = None
    col_hoa_mgt = None
    col_repairs_other = None
    col_mortgage_interest = None
    col_mortgage_principal = None
    col_apartment_improve = None
    col_owner_distribution = None
    col_total_outflow = None
    col_monthly = None
    col_cumulative = None

    for c in range(len(tbl.columns)):
        header = _norm_header(tbl.cell(1, c).text)
        header_compact = _hdr_compact(tbl.cell(1, c).text)

        if header == "Month Year":
            col_month_year = c
        elif header_compact == "Owner Contribution":
            col_owner_contrib = c
        elif header_compact == "Mortgage Loan":
            col_mortgage_loan = c
        elif header_compact == "Rent & Dividend":
            col_rent_dividend = c
        elif header_compact == "Total Inflow":
            col_total_inflow = c
        elif header_compact == "HOA & Mgt. Fee":
            col_hoa_mgt = c
        elif header_compact == "Repairs & Other Exp.":
            col_repairs_other = c
        elif header_compact == "Mortgage Interest":
            col_mortgage_interest = c
        elif header_compact in ("Mortgage Principal", "Mortgage\nPrincipal"):
            col_mortgage_principal = c
        elif header_compact in ("Apartment & Improve.", "Apartment & Improve"):
            col_apartment_improve = c
        elif header_compact == "Owner Distribution":
            col_owner_distribution = c
        elif header_compact == "Total Outflow":
            col_total_outflow = c
        elif header_compact == "Monthly":
            col_monthly = c
        elif header_compact == "Cumulative":
            col_cumulative = c

    if col_month_year is None:
        print("monthly_cash_table missing required column header: Month Year")
        return

    token_to_label = build_month_year_labels(ctx, property_name=None)

    sql = f"""
            SELECT timeframe,
                cash_categorization,
                cash_type_mapping,
                SUM(cash_value) AS total_value
            FROM gl_agg
            WHERE investor = ?
            AND (timeframe IS NULL OR timeframe <> 'N/A')
            AND timeframe IN ('[T1]','[T2]','[T3]','[T4]','[T5]','[T6]','[T7]','[T8]','[T9]','[T10]','[T11]','[T12]','[T13]')
            {owner_sql}
            GROUP BY timeframe, cash_categorization, cash_type_mapping
        """

    con = sqlite3.connect(str(config.SQLITE_PATH))
    rows = con.execute(sql, (ctx.investor, *owner_params)).fetchall()
    con.close()

    vals: Dict[str, Dict[str, float]] = {}
    totals_by_tf: Dict[str, Dict[str, float]] = {}

    for tf, cat, cash_type, total_value in rows:
        tf_key = str(tf).strip()
        cat_key = str(cat or "").strip()
        ct = str(cash_type or "").strip().lower()
        v_raw = -1.0 * float(total_value or 0.0)

        # Special handling: Mortgage Principal can be "Both" and must land in either
        # "Mortgage Loan" (if positive) or "Mortgage Principal" (if negative).
        if ct == "both" and cat_key == "Mortgage Principal":
            if v_raw > 0:
                cat_key = "Mortgage Loan"
                ct_eff = "inflow"
            elif v_raw < 0:
                cat_key = "Mortgage Principal"
                ct_eff = "outflow"
            else:
                ct_eff = ""
            v_signed = v_raw
        else:
            ct_eff = ct
            v_signed = 0.0
            if ct_eff in ("inflow", "outflow"):
                v_signed = v_raw
            else:
                v_signed = 0.0

        if tf_key not in vals:
            vals[tf_key] = {}
        if cat_key != "":
            vals[tf_key][cat_key] = vals[tf_key].get(cat_key, 0.0) + v_signed

        if tf_key not in totals_by_tf:
            totals_by_tf[tf_key] = {"inflow": 0.0, "outflow": 0.0}

        if ct_eff == "inflow":
            totals_by_tf[tf_key]["inflow"] += v_raw
        elif ct_eff == "outflow":
            totals_by_tf[tf_key]["outflow"] += v_raw

    total_row_idx = len(tbl.rows) - 1
    cumulative_running = 0.0

    total_owner_contrib = 0.0
    total_mortgage_loan = 0.0
    total_rent_dividend = 0.0
    total_inflow_all = 0.0

    total_hoa_mgt = 0.0
    total_repairs_other = 0.0
    total_mortgage_interest = 0.0
    total_mortgage_principal = 0.0
    total_apartment_improve = 0.0
    total_owner_distribution = 0.0
    total_outflow_all = 0.0

    data_row_count = max(0, len(tbl.rows) - 3)
    print(f"monthly_cash_table Starting process for {data_row_count} rows.")

    current = 0
    for r in range(2, len(tbl.rows)):
        if r == total_row_idx:
            continue

        row_label = tbl.cell(r, col_month_year).text.strip()
        if row_label == "":
            continue

        tf_token = None
        for token in (f"[T{n}]" for n in range(1, 14)):
            if token in row_label:
                tf_token = token
                break

        if tf_token is None:
            continue

        if tf_token not in token_to_label:
            continue

        current += 1
        print(f"monthly_cash_table Currently on {current} of {data_row_count}")

        new_label = row_label.replace(tf_token, token_to_label[tf_token])
        _set_cell_text_preserve_cell_format(tbl.cell(r, col_month_year), new_label)
        p = tbl.cell(r, col_month_year).text_frame.paragraphs[0]
        if p.runs:
            r0 = p.runs[0]
            r0.font.name = "Lato"
            r0.font.size = Pt(10)
            r0.font.color.rgb = RGBColor(0, 0, 0)

        tf_vals = vals.get(tf_token, {})
        tf_totals = totals_by_tf.get(tf_token, {"inflow": 0.0, "outflow": 0.0})

        owner_contrib = apply_ownership_amount(ctx, float(tf_vals.get("Owner Contribution", 0.0)), "monthly_cash_table.owner_contribution")
        mortgage_loan = apply_ownership_amount(ctx, float(tf_vals.get("Mortgage Loan", 0.0)), "monthly_cash_table.mortgage_loan")
        rent_dividend = apply_ownership_amount(ctx, float(tf_vals.get("Rent & Dividend", 0.0)), "monthly_cash_table.rent_dividend")

        hoa_mgt = apply_ownership_amount(ctx, float(tf_vals.get("HOA & Mgt. Fee", 0.0)), "monthly_cash_table.hoa_mgt_fee")
        repairs_other = apply_ownership_amount(ctx, float(tf_vals.get("Repairs & Other Exp.", 0.0)), "monthly_cash_table.repairs_other")
        mortgage_interest = apply_ownership_amount(ctx, float(tf_vals.get("Mortgage Interest", 0.0)), "monthly_cash_table.mortgage_interest")
        mortgage_principal = apply_ownership_amount(
            ctx,
            float(tf_vals.get("Mortgage Principal", tf_vals.get("Mortgage Principle", 0.0))),
            "monthly_cash_table.mortgage_principal",
        )
        apartment_improve = apply_ownership_amount(ctx, float(tf_vals.get("Apartment & Improve.", 0.0)), "monthly_cash_table.apartment_improve")
        owner_distribution = apply_ownership_amount(ctx, float(tf_vals.get("Owner Distribution", 0.0)), "monthly_cash_table.owner_distribution")

        total_inflow = apply_ownership_amount(ctx, float(tf_totals.get("inflow", 0.0)), "monthly_cash_table.total_inflow")
        total_outflow = apply_ownership_amount(ctx, float(tf_totals.get("outflow", 0.0)), "monthly_cash_table.total_outflow")

        monthly = total_inflow + total_outflow

        if cumulative_running == 0.0:
            cumulative_running = monthly
        else:
            cumulative_running = cumulative_running + monthly

        if col_owner_contrib is not None:
            _set_currency_cell(tbl.cell(r, col_owner_contrib), owner_contrib)
        if col_mortgage_loan is not None:
            _set_currency_cell(tbl.cell(r, col_mortgage_loan), mortgage_loan)
        if col_rent_dividend is not None:
            _set_currency_cell(tbl.cell(r, col_rent_dividend), rent_dividend)
        if col_total_inflow is not None:
            _set_currency_cell(tbl.cell(r, col_total_inflow), total_inflow)

        if col_hoa_mgt is not None:
            _set_currency_cell(tbl.cell(r, col_hoa_mgt), hoa_mgt)
        if col_repairs_other is not None:
            _set_currency_cell(tbl.cell(r, col_repairs_other), repairs_other)
        if col_mortgage_interest is not None:
            _set_currency_cell(tbl.cell(r, col_mortgage_interest), mortgage_interest)
        if col_mortgage_principal is not None:
            _set_currency_cell(tbl.cell(r, col_mortgage_principal), mortgage_principal)
        if col_apartment_improve is not None:
            _set_currency_cell(tbl.cell(r, col_apartment_improve), apartment_improve)
        if col_owner_distribution is not None:
            _set_currency_cell(tbl.cell(r, col_owner_distribution), owner_distribution)
        if col_total_outflow is not None:
            _set_currency_cell(tbl.cell(r, col_total_outflow), total_outflow)

        if col_monthly is not None:
            _set_currency_cell(tbl.cell(r, col_monthly), monthly)
        if col_cumulative is not None:
            _set_currency_cell(tbl.cell(r, col_cumulative), cumulative_running)

        total_owner_contrib += owner_contrib
        total_mortgage_loan += mortgage_loan
        total_rent_dividend += rent_dividend
        total_inflow_all += total_inflow

        total_hoa_mgt += hoa_mgt
        total_repairs_other += repairs_other
        total_mortgage_interest += mortgage_interest
        total_mortgage_principal += mortgage_principal
        total_apartment_improve += apartment_improve
        total_owner_distribution += owner_distribution
        total_outflow_all += total_outflow

    if col_owner_contrib is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_owner_contrib), total_owner_contrib)
    if col_mortgage_loan is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_mortgage_loan), total_mortgage_loan)
    if col_rent_dividend is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_rent_dividend), total_rent_dividend)
    if col_total_inflow is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_total_inflow), total_inflow_all)

    if col_hoa_mgt is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_hoa_mgt), total_hoa_mgt)
    if col_repairs_other is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_repairs_other), total_repairs_other)
    if col_mortgage_interest is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_mortgage_interest), total_mortgage_interest)
    if col_mortgage_principal is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_mortgage_principal), total_mortgage_principal)
    if col_apartment_improve is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_apartment_improve), total_apartment_improve)
    if col_owner_distribution is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_owner_distribution), total_owner_distribution)
    if col_total_outflow is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_total_outflow), total_outflow_all)

    total_monthly_all = total_inflow_all + total_outflow_all

    if col_monthly is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_monthly), total_monthly_all)
    if col_cumulative is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_cumulative), cumulative_running)

    print("monthly_cash_table updated.")

def update_available_cash(slide: Slide, shape: BaseShape, prs: Presentation, ctx: UpdateContext) -> None:
    if not hasattr(shape, "table"):
        return

    owner_sql, owner_params = _owner_filter_sql(ctx)

    from pptx.dml.color import RGBColor

    def _norm_header(s: str) -> str:
        return (s or "").replace("\r", "").replace(" \n", "\n").strip()

    def _fmt_currency(x: float) -> tuple[str, bool]:
        if abs(x) < 0.5:
            return "-", False
        if x < 0:
            return f"(${abs(x):,.0f})", True
        return f"${x:,.0f}", False

    def _apply_red_if_negative(cell, is_negative: bool) -> None:
        if not is_negative:
            return
        try:
            p0 = cell.text_frame.paragraphs[0]
            if p0.runs:
                p0.runs[0].font.color.rgb = RGBColor(255, 0, 0)
        except Exception:
            return

    def _set_currency_cell(cell, amount: float) -> None:
        txt, is_neg = _fmt_currency(amount)
        _set_cell_text_preserve_cell_format(cell, txt)
        _apply_red_if_negative(cell, is_neg)

    tbl = shape.table

    # Locate the header row. Some PPT tables use row 0, others use row 1 (like your other tables).
    header_row_idx = None
    for candidate in (0, 1):
        if candidate >= len(tbl.rows):
            continue
        row_headers = [_norm_header(tbl.cell(candidate, c).text) for c in range(len(tbl.columns))]
        if (
            "Reserve\nAccount Balance" in row_headers
            or "Reserve Account Balance" in row_headers
            or "Investor\nAccount Balance" in row_headers
            or "Investor Account Balance" in row_headers
            or "Current\nAvailable Cash" in row_headers
            or "Current Available Cash" in row_headers
        ):
            header_row_idx = candidate
            break

    if header_row_idx is None:
        print("available_cash missing required header row")
        return

    value_row_idx = header_row_idx + 1
    if value_row_idx >= len(tbl.rows):
        print("available_cash missing value row beneath header row")
        return

    col_reserve = None
    col_investor = None
    col_available = None

    for c in range(len(tbl.columns)):
        header = _norm_header(tbl.cell(header_row_idx, c).text)
        if header in ("Reserve\nAccount Balance", "Reserve Account Balance"):
            col_reserve = c
        elif header in ("Investor\nAccount Balance", "Investor Account Balance"):
            col_investor = c
        elif header in ("Current\nAvailable Cash", "Current Available Cash"):
            col_available = c

    if col_reserve is None or col_investor is None or col_available is None:
        print("available_cash missing one or more required column headers")
        return

    # DB rule you provided: cash_value negative for inflow, positive for outflow.
    # For account balances, we present positive cash as a positive number, so we flip sign.
    sql = f"""
            SELECT
                COALESCE(SUM(CASE WHEN cash_categorization = '1180 Cash Account' THEN cash_value ELSE 0 END), 0.0) AS reserve_balance,
                COALESCE(SUM(CASE WHEN cash_categorization = '1150 Cash Account' THEN cash_value ELSE 0 END), 0.0) AS investor_balance
            FROM gl_agg
            WHERE investor = ?
            AND (timeframe IS NULL OR timeframe <> 'N/A')
            {owner_sql}
        """

    con = sqlite3.connect(str(config.SQLITE_PATH))
    row = con.execute(sql, (ctx.investor, *owner_params)).fetchone()
    con.close()

    reserve_raw = float((row[0] if row and row[0] is not None else 0.0))
    investor_raw = float((row[1] if row and row[1] is not None else 0.0))

    reserve_balance = apply_ownership_amount(ctx, reserve_raw, "available_cash.reserve_balance")
    investor_balance = apply_ownership_amount(ctx, investor_raw, "available_cash.investor_balance")
    current_available = apply_ownership_amount(ctx, reserve_raw + investor_raw, "available_cash.current_available_cash")

    _set_currency_cell(tbl.cell(value_row_idx, col_reserve), reserve_balance)
    _set_currency_cell(tbl.cell(value_row_idx, col_investor), investor_balance)
    _set_currency_cell(tbl.cell(value_row_idx, col_available), current_available)

    print("available_cash updated.")
