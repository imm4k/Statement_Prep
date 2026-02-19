from __future__ import annotations

import sqlite3
from typing import Dict

import config

from pptx.presentation import Presentation
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from ppt_objects import UpdateContext
from ppt_text_replace import replace_tokens_in_shape

def _owner_filter_sql(ctx: UpdateContext) -> tuple[str, tuple]:
    if ctx.owner is None or str(ctx.owner).strip() == "":
        return "", tuple()
    return " AND owner = ? ", (str(ctx.owner).strip(),)


def _replace_tokens_in_shape_robust(shape: BaseShape, token_map: Dict[str, str]) -> int:
    """
    Fixes tokens split across multiple runs (your [Cumulative Return] case).
    Strategy:
      1) Try existing replace_tokens_in_shape
      2) Then do an across runs pass per paragraph by rewriting run0 and blanking the rest
    """
    count = 0

    try:
        count += int(replace_tokens_in_shape(shape, token_map) or 0)
    except Exception:
        pass

    if not getattr(shape, "has_text_frame", False):
        return count

    tf = shape.text_frame
    for p in tf.paragraphs:
        runs = list(p.runs)
        if not runs:
            continue

        full = "".join(r.text for r in runs)
        new = full

        para_repls = 0
        for k, v in token_map.items():
            if not k:
                continue
            occurrences = new.count(k)
            if occurrences:
                para_repls += occurrences
                new = new.replace(k, v)

        if new != full:
            runs[0].text = new
            for r in runs[1:]:
                r.text = ""
            count += para_repls

    return count


def update_cover_title(slide: Slide, shape: BaseShape, prs: Presentation, ctx: UpdateContext) -> None:
    token_map = {"[T1]": ctx.t1_str}
    count = _replace_tokens_in_shape_robust(shape, token_map)
    print(f"cover_title replacements applied: {count}")


def update_cover_subtitle(slide: Slide, shape: BaseShape, prs: Presentation, ctx: UpdateContext) -> None:
    token_map = {"[Investor]": ctx.investor}
    count = _replace_tokens_in_shape_robust(shape, token_map)
    print(f"cover_subtitle replacements applied: {count}")


def _join_owner_list_for_display(items) -> str:
    owners = [str(x).strip() for x in (items or []) if x is not None and str(x).strip() != ""]
    owners = sorted(list(set(owners)))
    if not owners:
        return ""
    if len(owners) == 1:
        return owners[0]
    if len(owners) == 2:
        return f"{owners[0]} and {owners[1]}"
    return f"{', '.join(owners[:-1])}, and {owners[-1]}"


def _fmt_percent_1dp(x: float) -> str:
    return f"{x * 100.0:.1f}%"


def _fmt_usd_short_k(x: float) -> str:
    k = int(round(float(x or 0.0) / 1000.0))
    return f"${k:,}K"


def _get_investor_owners(investor: str) -> str:
    sql = """
        SELECT DISTINCT owner
        FROM gl_agg
        WHERE investor = ?
          AND owner IS NOT NULL
          AND TRIM(owner) <> ''
          AND (timeframe IS NULL OR timeframe <> 'N/A')
        ORDER BY owner
    """
    con = sqlite3.connect(str(config.SQLITE_PATH))
    rows = con.execute(sql, (investor,)).fetchall()
    con.close()
    return _join_owner_list_for_display([r[0] for r in rows])

def _get_portfolio_total_invested(ctx: UpdateContext) -> float:
    owner_sql, owner_params = _owner_filter_sql(ctx)
    sql = f"""
        SELECT ABS(SUM(value)) AS total_invested
        FROM gl_agg
        WHERE investor = ?
          AND categorization = 'Total Invested'
          AND (timeframe IS NULL OR timeframe <> 'N/A')
          {owner_sql}
    """
    con = sqlite3.connect(str(config.SQLITE_PATH))
    row = con.execute(sql, (ctx.investor, *owner_params)).fetchone()
    con.close()
    return float((row[0] if row and row[0] is not None else 0.0))

def _get_portfolio_cumulative_return_amount(ctx: UpdateContext) -> float:
    owner_sql, owner_params = _owner_filter_sql(ctx)
    sql = f"""
        SELECT
            ABS(SUM(CASE WHEN categorization = 'Total Invested' THEN value ELSE 0 END)) AS invested,
            ABS(SUM(CASE WHEN categorization = 'Mortgage Balance' THEN value ELSE 0 END)) AS mortgage,
            SUM(
                CASE
                    WHEN UPPER(TRIM(COALESCE(gl_mapping_type, ''))) = 'REVENUE' THEN -1.0 * value
                    WHEN UPPER(TRIM(COALESCE(gl_mapping_type, ''))) = 'EXPENSE' THEN -1.0 * value
                    ELSE 0.0
                END
            ) AS income
        FROM gl_agg
        WHERE investor = ?
          AND (timeframe IS NULL OR timeframe <> 'N/A')
          {owner_sql}
    """
    con = sqlite3.connect(str(config.SQLITE_PATH))
    row = con.execute(sql, (ctx.investor, *owner_params)).fetchone()
    con.close()

    invested = float(row[0] or 0.0)
    mortgage = float(row[1] or 0.0)
    income = float(row[2] or 0.0)

    nav = -mortgage
    return nav + income - invested

def _get_portfolio_cumulative_income(ctx: UpdateContext) -> float:
    owner_sql, owner_params = _owner_filter_sql(ctx)
    sql = f"""
        SELECT
            SUM(
                CASE
                    WHEN UPPER(TRIM(COALESCE(gl_mapping_type, ''))) = 'REVENUE' THEN -1.0 * value
                    WHEN UPPER(TRIM(COALESCE(gl_mapping_type, ''))) = 'EXPENSE' THEN -1.0 * value
                    ELSE 0.0
                END
            ) AS cumulative_income
        FROM gl_agg
        WHERE investor = ?
          AND (timeframe IS NULL OR timeframe <> 'N/A')
          AND timeframe IN ('[T1]','[T2]','[T3]','[T4]','[T5]','[T6]','[T7]','[T8]','[T9]','[T10]','[T11]','[T12]','[T13]')
          {owner_sql}
    """
    con = sqlite3.connect(str(config.SQLITE_PATH))
    row = con.execute(sql, (ctx.investor, *owner_params)).fetchone()
    con.close()
    return float((row[0] if row and row[0] is not None else 0.0))


def update_summary_title(slide: Slide, shape: BaseShape, prs: Presentation, ctx: UpdateContext) -> None:
    owner_str = str(ctx.owner).strip() if ctx.owner else _get_investor_owners(ctx.investor)
    token_map = {
        "[Owner]": owner_str,
        "[T1]": ctx.t1_str,
    }
    count = _replace_tokens_in_shape_robust(shape, token_map)
    print(f"summary_title replacements applied: {count}")

def update_cash_summary_title(slide: Slide, shape: BaseShape, prs: Presentation, ctx: UpdateContext) -> None:
    owner_str = str(ctx.owner).strip() if ctx.owner else _get_investor_owners(ctx.investor)
    token_map = {
        "[Owner]": owner_str,
        "[T1]": ctx.t1_str,
    }
    count = _replace_tokens_in_shape_robust(shape, token_map)
    print(f"cash_summary_title replacements applied: {count}")

def update_summary_top_text(slide: Slide, shape: BaseShape, prs: Presentation, ctx: UpdateContext) -> None:
    def _norm_header(s: str) -> str:
        return (s or "").replace("\r", "").replace(" \n", "\n").strip()

    def _parse_currency(s: str) -> float:
        t = (s or "").strip()
        if t == "" or t == "-":
            return 0.0
        neg = False
        if t.startswith("(") and t.endswith(")"):
            neg = True
            t = t[1:-1]
        t = t.replace("$", "").replace(",", "").strip()
        try:
            v = float(t)
        except Exception:
            v = 0.0
        return -v if neg else v

    def _parse_percent(s: str) -> float:
        t = (s or "").strip()
        if t == "" or t == "-":
            return 0.0
        t = t.replace("%", "").replace(",", "").strip()
        try:
            return float(t) / 100.0
        except Exception:
            return 0.0

    owner_str = str(ctx.owner).strip() if ctx.owner else _get_investor_owners(ctx.investor)

    total_invested = 0.0
    pct_return = 0.0

    summary_tbl = None
    for shp in slide.shapes:
        if getattr(shp, "name", "") == "summary_table" and hasattr(shp, "table"):
            summary_tbl = shp.table
            break

    if summary_tbl is None:
        print("summary_top_text: summary_table not found on this slide")
    else:
        total_row_idx = len(summary_tbl.rows) - 1

        col_total_invested = None
        col_pct_return = None

        for c in range(len(summary_tbl.columns)):
            header = _norm_header(summary_tbl.cell(1, c).text)
            if header == "Total\nInvested":
                col_total_invested = c
            elif header == "% Return":
                col_pct_return = c

        if col_total_invested is None or col_pct_return is None:
            print("summary_top_text: required columns not found in summary_table")
        else:
            total_invested = abs(_parse_currency(summary_tbl.cell(total_row_idx, col_total_invested).text))
            pct_return = _parse_percent(summary_tbl.cell(total_row_idx, col_pct_return).text)

    cumulative_income = _get_portfolio_cumulative_income(ctx)

    coc = 0.0
    if abs(total_invested) > 0.0000001:
        coc = cumulative_income / total_invested

    token_map = {
        "[Owner]": owner_str,
        "[CoC Return]": _fmt_percent_1dp(coc),
        "[Cumulative Return]": _fmt_percent_1dp(pct_return),
        "[Total Invested (Short)]": _fmt_usd_short_k(total_invested),
    }

    count = _replace_tokens_in_shape_robust(shape, token_map)
    print(f"summary_top_text replacements applied: {count}")