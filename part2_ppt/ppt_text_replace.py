from __future__ import annotations

from typing import Dict

from pptx.shapes.base import BaseShape


def replace_tokens_in_shape(shape: BaseShape, token_map: Dict[str, str]) -> int:
    if not hasattr(shape, "text_frame"):
        return 0

    tf = shape.text_frame
    replacements = 0

    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            if run.text is None:
                continue
            new_text = run.text
            for token, value in token_map.items():
                if token in new_text:
                    new_text = new_text.replace(token, value)
            if new_text != run.text:
                run.text = new_text
                replacements += 1

    return replacements
