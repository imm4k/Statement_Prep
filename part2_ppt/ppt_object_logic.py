from __future__ import annotations

import sqlite3
from datetime import datetime
from typing import Dict

import config

from pptx.presentation import Presentation
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from lxml import etree

from ppt_objects import ObjectUpdater, UpdateContext
from ppt_text_replace import replace_tokens_in_shape


def _update_cover_title(slide: Slide, shape: BaseShape, prs: Presentation, ctx: UpdateContext) -> None:
    token_map = {"[T1]": ctx.t1_str}
    count = replace_tokens_in_shape(shape, token_map)
    print(f"cover_title replacements applied: {count}")


def _update_cover_subtitle(slide: Slide, shape: BaseShape, prs: Presentation, ctx: UpdateContext) -> None:
    token_map = {"[Investor]": ctx.investor}
    count = replace_tokens_in_shape(shape, token_map)
    print(f"cover_subtitle replacements applied: {count}")

def _set_cell_text_preserve_cell_format(cell, text: str) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    txBody = cell._tc.txBody
    p = txBody.find("a:p", ns)
    if p is None:
        cell.text_frame.text = text
        return

    pPr = p.find("a:pPr", ns)
    endParaRPr = p.find("a:endParaRPr", ns)

    algn = None
    if pPr is not None:
        algn = pPr.get("algn")

    typeface = None
    sz = None
    b = None
    i = None
    u = None
    color_val = None

    if endParaRPr is not None:
        sz = endParaRPr.get("sz")
        b = endParaRPr.get("b")
        i = endParaRPr.get("i")
        u = endParaRPr.get("u")

        latin = endParaRPr.find("a:latin", ns)
        if latin is not None:
            typeface = latin.get("typeface")

        srgb = endParaRPr.find(".//a:srgbClr", ns)
        if srgb is not None:
            color_val = srgb.get("val")

    cell.text_frame.text = text

    p0 = cell.text_frame.paragraphs[0]
    if algn == "ctr":
        p0.alignment = PP_ALIGN.CENTER
    elif algn == "l":
        p0.alignment = PP_ALIGN.LEFT
    elif algn == "r":
        p0.alignment = PP_ALIGN.RIGHT
    elif algn == "just":
        p0.alignment = PP_ALIGN.JUSTIFY

    if p0.runs:
        r0 = p0.runs[0]
        if typeface:
            r0.font.name = typeface
        if sz and str(sz).isdigit():
            r0.font.size = Pt(int(sz) / 100)
        if b is not None:
            r0.font.bold = (str(b) == "1")
        if i is not None:
            r0.font.italic = (str(i) == "1")
        if u is not None:
            r0.font.underline = (str(u).lower() != "none")
        if color_val and len(color_val) == 6:
            r0.font.color.rgb = RGBColor(int(color_val[0:2], 16), int(color_val[2:4], 16), int(color_val[4:6], 16))

def _update_summary_table(slide: Slide, shape: BaseShape, prs: Presentation, ctx: UpdateContext) -> None:
    if not hasattr(shape, "table"):
        return

    from openpyxl import load_workbook
    from pptx.dml.color import RGBColor

    def _norm_header(s: str) -> str:
        return s.replace("\r", "").replace(" \n", "\n").strip()

    def _fmt_currency(x: float) -> tuple[str, bool]:
        # Matches Excel custom format: $#,##0;[Red]($#,##0);-
        if abs(x) < 0.5:
            return "-", False
        if x < 0:
            return f"(${abs(x):,.0f})", True
        return f"${x:,.0f}", False

    def _apply_red_if_negative(cell, is_negative: bool) -> None:
        if not is_negative:
            return
        try:
            p0 = cell.text_frame.paragraphs[0]
            if p0.runs:
                p0.runs[0].font.color.rgb = RGBColor(255, 0, 0)
        except Exception:
            return

    def _set_currency_cell(cell, amount: float) -> None:
        txt, is_neg = _fmt_currency(amount)
        _set_cell_text_preserve_cell_format(cell, txt)
        _apply_red_if_negative(cell, is_neg)

    def _fmt_percent(x: float) -> tuple[str, bool]:
        # Format as 0.0%
        val = x * 100.0
        txt = f"{val:.1f}%"
        return txt, (x < 0)

    def _set_percent_cell(cell, pct: float) -> None:
        txt, is_neg = _fmt_percent(pct)
        _set_cell_text_preserve_cell_format(cell, txt)
        _apply_red_if_negative(cell, is_neg)

    def _read_general_config_market_values() -> Dict[str, float]:
        labels = {
            "Studio": "Studio Market:",
            "1-Bed": "1-Bed Market:",
            "2-Bed": "2-Bed Market:",
            "3-Bed": "3-Bed Market:",
        }

        wb = load_workbook(filename=str(config.SETUP_EXCEL_PATH), data_only=True)
        try:
            if config.GENERAL_CONFIG_SHEET not in wb.sheetnames:
                raise RuntimeError(f"Missing sheet: {config.GENERAL_CONFIG_SHEET}")

            ws = wb[config.GENERAL_CONFIG_SHEET]
            label_to_value: Dict[str, float] = {}

            # Column A labels, column B values
            for r in range(1, ws.max_row + 1):
                a = ws.cell(r, 1).value
                b = ws.cell(r, 2).value
                if a is None:
                    continue
                key = str(a).strip()
                if key in labels.values():
                    try:
                        label_to_value[key] = float(b)
                    except Exception:
                        label_to_value[key] = 0.0

            out: Dict[str, float] = {}
            for unit_type, label in labels.items():
                out[unit_type] = float(label_to_value.get(label, 0.0))

            return out
        finally:
            wb.close()

    tbl = shape.table

    # Column headers on row index 1 (PowerPoint row 2)
    col_property = None
    col_type = None
    col_duration = None
    col_total_invested = None
    col_est_mkt_value = None
    col_mortgage_balance = None
    col_nav = None
    col_cum_income = None
    col_cum_return = None
    col_pct_return = None

    for c in range(len(tbl.columns)):
        header = _norm_header(tbl.cell(1, c).text)

        if header == "Property":
            col_property = c
        elif header == "Type":
            col_type = c
        elif header == "Duration\n(Months)":
            col_duration = c
        elif header == "Total\nInvested":
            col_total_invested = c
        elif header == "Estimated\nMarket Value":
            col_est_mkt_value = c
        elif header == "Mortgage\nBalance":
            col_mortgage_balance = c
        elif header == "Net Asset Value (NAV)":
            col_nav = c
        elif header == "Cumulative Income":
            col_cum_income = c
        elif header == "Cumulative Return":
            col_cum_return = c
        elif header == "% Return":
            col_pct_return = c

    if col_property is None:
        print("summary_table missing required column header: Property")
        return

    total_row_idx = len(tbl.rows) - 1

    market_values_by_type = _read_general_config_market_values()

    con = sqlite3.connect(str(config.SQLITE_PATH))

    sql_durations = """
        SELECT property, MIN(acquired)
        FROM gl_agg
        WHERE investor = ?
          AND acquired IS NOT NULL
        GROUP BY property
    """

    sql_total_invested = """
        SELECT property,
               ABS(SUM(value)) AS total_invested
        FROM gl_agg
        WHERE investor = ?
          AND categorization = 'Total Invested'
          AND (timeframe IS NULL OR timeframe <> 'N/A')
        GROUP BY property
    """

    sql_mortgage_balance = """
        SELECT property,
               ABS(SUM(value)) AS mortgage_balance
        FROM gl_agg
        WHERE investor = ?
          AND categorization = 'Mortgage Balance'
          AND (timeframe IS NULL OR timeframe <> 'N/A')
        GROUP BY property
    """

    sql_income = """
        SELECT property,
               SUM(
                   CASE
                       WHEN UPPER(TRIM(COALESCE(gl_mapping_type, ''))) = 'REVENUE' THEN -1.0 * value
                       WHEN UPPER(TRIM(COALESCE(gl_mapping_type, ''))) = 'EXPENSE' THEN -1.0 * value
                       ELSE 0.0
                   END
               ) AS cumulative_income
        FROM gl_agg
        WHERE investor = ?
          AND (timeframe IS NULL OR timeframe <> 'N/A')
        GROUP BY property
    """

    duration_rows = con.execute(sql_durations, (ctx.investor,)).fetchall()
    invested_rows = con.execute(sql_total_invested, (ctx.investor,)).fetchall()
    mortgage_rows = con.execute(sql_mortgage_balance, (ctx.investor,)).fetchall()
    income_rows = con.execute(sql_income, (ctx.investor,)).fetchall()

    con.close()

    durations: Dict[str, float] = {}
    for prop, acquired_val in duration_rows:
        acquired_dt = datetime.strptime(str(acquired_val)[:10], "%Y-%m-%d")
        months = (
            (ctx.statement_thru_date_dt.year - acquired_dt.year) * 12
            + (ctx.statement_thru_date_dt.month - acquired_dt.month)
        )
        if ctx.statement_thru_date_dt.day < acquired_dt.day:
            months -= 1
        durations[str(prop)] = float(months)

    total_invested_by_prop: Dict[str, float] = {}
    for prop, v in invested_rows:
        if prop is None:
            continue
        total_invested_by_prop[str(prop)] = float(v or 0.0)

    mortgage_by_prop: Dict[str, float] = {}
    for prop, v in mortgage_rows:
        if prop is None:
            continue
        mortgage_by_prop[str(prop)] = float(v or 0.0)

    income_by_prop: Dict[str, float] = {}
    for prop, v in income_rows:
        if prop is None:
            continue
        income_by_prop[str(prop)] = float(v or 0.0)

    duration_hits = 0
    invested_hits = 0
    est_hits = 0
    mortgage_hits = 0
    nav_hits = 0
    income_hits = 0
    return_hits = 0

    invested_total = 0.0
    est_total = 0.0
    mortgage_total = 0.0
    nav_total = 0.0
    income_total = 0.0
    cum_return_total = 0.0

    data_row_count = max(0, len(tbl.rows) - 3)
    print(f"summary_table Starting process for {data_row_count} rows.")

    # Data rows start at index 2, final row is Total
    current = 0
    for r in range(2, len(tbl.rows)):
        if r == total_row_idx:
            continue

        current += 1
        print(f"summary_table Currently on {current} of {data_row_count}")

        prop_name = tbl.cell(r, col_property).text.strip()
        if prop_name == "":
            continue

        if col_duration is not None and prop_name in durations:
            _set_cell_text_preserve_cell_format(tbl.cell(r, col_duration), str(durations[prop_name]))
            duration_hits += 1

        total_invested = abs(total_invested_by_prop.get(prop_name, 0.0))
        mortgage_bal = abs(mortgage_by_prop.get(prop_name, 0.0))
        cum_income = float(income_by_prop.get(prop_name, 0.0))

        unit_type = ""
        if col_type is not None:
            unit_type = tbl.cell(r, col_type).text.strip()

        if unit_type == "" and col_type is not None:
            continue

        est_mkt = float(market_values_by_type.get(unit_type, 0.0))

        nav = est_mkt - mortgage_bal
        cum_return = nav + cum_income - total_invested

        if col_total_invested is not None:
            _set_currency_cell(tbl.cell(r, col_total_invested), total_invested)
            invested_hits += 1

        if col_est_mkt_value is not None:
            _set_currency_cell(tbl.cell(r, col_est_mkt_value), est_mkt)
            est_hits += 1

        if col_mortgage_balance is not None:
            _set_currency_cell(tbl.cell(r, col_mortgage_balance), mortgage_bal)
            mortgage_hits += 1

        if col_nav is not None:
            _set_currency_cell(tbl.cell(r, col_nav), nav)
            nav_hits += 1

        if col_cum_income is not None:
            _set_currency_cell(tbl.cell(r, col_cum_income), cum_income)
            income_hits += 1

        if col_cum_return is not None:
            _set_currency_cell(tbl.cell(r, col_cum_return), cum_return)
            return_hits += 1

        invested_total += total_invested
        est_total += est_mkt
        mortgage_total += mortgage_bal
        nav_total += nav
        income_total += cum_income
        cum_return_total += cum_return

    # Total row updates
    if col_total_invested is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_total_invested), invested_total)

    if col_est_mkt_value is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_est_mkt_value), est_total)

    if col_mortgage_balance is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_mortgage_balance), mortgage_total)

    if col_nav is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_nav), nav_total)

    if col_cum_income is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_cum_income), income_total)

    if col_cum_return is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_cum_return), cum_return_total)

    if col_pct_return is not None:
        pct = 0.0
        if abs(invested_total) > 0.0000001:
            pct = cum_return_total / invested_total
        _set_percent_cell(tbl.cell(total_row_idx, col_pct_return), pct)

    print(f"summary_table duration updated rows: {duration_hits}")
    print(f"summary_table total_invested updated rows: {invested_hits}")
    print(f"summary_table estimated_market_value updated rows: {est_hits}")
    print(f"summary_table mortgage_balance updated rows: {mortgage_hits}")
    print(f"summary_table nav updated rows: {nav_hits}")
    print(f"summary_table cumulative_income updated rows: {income_hits}")
    print(f"summary_table cumulative_return updated rows: {return_hits}")

def _update_monthly_perf_table(slide: Slide, shape: BaseShape, prs: Presentation, ctx: UpdateContext) -> None:
    if not hasattr(shape, "table"):
        return

    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    from datetime import date

    def _norm_header(s: str) -> str:
        return s.replace("\r", "").replace(" \n", "\n").strip()

    def _fmt_currency(x: float) -> tuple[str, bool]:
        # Matches Excel custom format: $#,##0;[Red]($#,##0);-
        if abs(x) < 0.5:
            return "-", False
        if x < 0:
            return f"(${abs(x):,.0f})", True
        return f"${x:,.0f}", False

    def _apply_red_if_negative(cell, is_negative: bool) -> None:
        if not is_negative:
            return
        try:
            p0 = cell.text_frame.paragraphs[0]
            if p0.runs:
                p0.runs[0].font.color.rgb = RGBColor(255, 0, 0)
        except Exception:
            return

    def _set_currency_cell(cell, amount: float) -> None:
        txt, is_neg = _fmt_currency(amount)
        _set_cell_text_preserve_cell_format(cell, txt)
        _apply_red_if_negative(cell, is_neg)

    def _month_year_label(dt: date) -> str:
        # Example: "Mar 2024"
        return dt.strftime("%b %Y")

    def _add_months(dt: date, months: int) -> date:
        y = dt.year + (dt.month - 1 + months) // 12
        m = (dt.month - 1 + months) % 12 + 1
        return date(y, m, 1)

    tbl = shape.table

    # Column headers on row index 1 (PowerPoint row 2)
    col_month_year = None
    col_rent = None
    col_other_rev = None
    col_total_rev = None
    col_mortgage = None
    col_hoa = None
    col_mgt_fee = None
    col_repairs = None
    col_other_exp = None
    col_total_exp = None
    col_monthly = None
    col_cumulative = None

    for c in range(len(tbl.columns)):
        header = _norm_header(tbl.cell(1, c).text)
        if header == "Month Year":
            col_month_year = c
        elif header == "Rent":
            col_rent = c
        elif header == "Other Revenue":
            col_other_rev = c
        elif header == "Total Revenue":
            col_total_rev = c
        elif header == "Mortgage":
            col_mortgage = c
        elif header == "HOA":
            col_hoa = c
        elif header == "Mgt. Fee":
            col_mgt_fee = c
        elif header == "Repairs Exp.":
            col_repairs = c
        elif header == "Other Expense":
            col_other_exp = c
        elif header == "Total Expenses":
            col_total_exp = c
        elif header == "Monthly":
            col_monthly = c
        elif header == "Cumulative":
            col_cumulative = c

    if col_month_year is None:
        print("monthly_perf_table missing required column header: Month Year")
        return

    # Build token map: [T13] is statement thru month, [T12] is prior month, ... [T1] is 12 months prior
    stmt_month_start = date(ctx.statement_thru_date_dt.year, ctx.statement_thru_date_dt.month, 1)

    token_to_month_start: Dict[str, date] = {}
    for n in range(1, 14):
        months_back = 13 - n  # T13 -> 0 back, T1 -> 12 back
        token_to_month_start[f"[T{n}]"] = _add_months(stmt_month_start, -months_back)

    token_to_label: Dict[str, str] = {k: _month_year_label(v) for k, v in token_to_month_start.items()}

    # Pull all needed monthly values in one query (excluding timeframe = 'N/A')
    wanted_cats = (
        "Rent",
        "Other Revenue",
        "Mortgage",
        "HOA",
        "Mgt. Fee",
        "Repairs Exp.",
        "Other Expense",
    )

    placeholders = ",".join(["?"] * len(wanted_cats))
    sql = f"""
        SELECT timeframe, categorization, gl_mapping_type, SUM(value) AS total_value
        FROM gl_agg
        WHERE investor = ?
          AND (timeframe IS NULL OR timeframe <> 'N/A')
          AND timeframe IN ('[T1]','[T2]','[T3]','[T4]','[T5]','[T6]','[T7]','[T8]','[T9]','[T10]','[T11]','[T12]','[T13]')
          AND categorization IN ({placeholders})
        GROUP BY timeframe, categorization, gl_mapping_type
    """

    con = sqlite3.connect(str(config.SQLITE_PATH))
    rows = con.execute(sql, (ctx.investor, *wanted_cats)).fetchall()
    con.close()

    # Build dict: timeframe -> categorization -> adjusted_value
    # Adjustment: for Revenue and Expense mapping types, flip sign (debit-credit convention)
    vals: Dict[str, Dict[str, float]] = {}
    for tf, cat, mapping_type, total_value in rows:
        tf_key = str(tf).strip()
        cat_key = str(cat).strip()
        mt = str(mapping_type or "").strip().lower()
        v = float(total_value or 0.0)

        if mt in ("revenue", "expense"):
            v = -1.0 * v

        if tf_key not in vals:
            vals[tf_key] = {}
        vals[tf_key][cat_key] = vals[tf_key].get(cat_key, 0.0) + v

    total_row_idx = len(tbl.rows) - 1

    cumulative_running = 0.0

    # For totals row calculations
    total_rent = 0.0
    total_other_rev = 0.0
    total_mortgage = 0.0
    total_hoa = 0.0
    total_mgt_fee = 0.0
    total_repairs = 0.0
    total_other_exp = 0.0
    total_monthly = 0.0

    data_row_count = max(0, len(tbl.rows) - 3)
    print(f"monthly_perf_table Starting process for {data_row_count} rows.")

    current = 0
    for r in range(2, len(tbl.rows)):
        if r == total_row_idx:
            continue

        row_label = tbl.cell(r, col_month_year).text.strip()
        if row_label == "":
            continue

        # Identify timeframe token present in the row label
        tf_token = None
        for token in token_to_label.keys():
            if token in row_label:
                tf_token = token
                break

        if tf_token is None:
            continue

        current += 1
        print(f"monthly_perf_table Currently on {current} of {data_row_count}")

        # Replace the token in the Month Year label cell
        new_label = row_label.replace(tf_token, token_to_label[tf_token])
        _set_cell_text_preserve_cell_format(tbl.cell(r, col_month_year), new_label)
        p = tbl.cell(r, col_month_year).text_frame.paragraphs[0]
        if p.runs:
            r0 = p.runs[0]
            r0.font.name = "Lato"
            r0.font.size = Pt(10)
            r0.font.color.rgb = RGBColor(0, 0, 0)

        tf_key = tf_token  # matches gl_agg timeframe strings like "[T13]"
        tf_vals = vals.get(tf_key, {})

        rent = float(tf_vals.get("Rent", 0.0))
        other_rev = float(tf_vals.get("Other Revenue", 0.0))

        mortgage = float(tf_vals.get("Mortgage", 0.0))
        hoa = float(tf_vals.get("HOA", 0.0))
        mgt_fee = float(tf_vals.get("Mgt. Fee", 0.0))
        repairs = float(tf_vals.get("Repairs Exp.", 0.0))
        other_exp = float(tf_vals.get("Other Expense", 0.0))

        total_rev = rent + other_rev
        total_exp = mortgage + hoa + mgt_fee + repairs + other_exp
        monthly = total_rev + total_exp

        # Cumulative is rolling down the table
        if cumulative_running == 0.0:
            cumulative_running = monthly
        else:
            cumulative_running = cumulative_running + monthly

        # Write base columns
        if col_rent is not None:
            _set_currency_cell(tbl.cell(r, col_rent), rent)
        if col_other_rev is not None:
            _set_currency_cell(tbl.cell(r, col_other_rev), other_rev)
        if col_total_rev is not None:
            _set_currency_cell(tbl.cell(r, col_total_rev), total_rev)

        if col_mortgage is not None:
            _set_currency_cell(tbl.cell(r, col_mortgage), mortgage)
        if col_hoa is not None:
            _set_currency_cell(tbl.cell(r, col_hoa), hoa)
        if col_mgt_fee is not None:
            _set_currency_cell(tbl.cell(r, col_mgt_fee), mgt_fee)
        if col_repairs is not None:
            _set_currency_cell(tbl.cell(r, col_repairs), repairs)
        if col_other_exp is not None:
            _set_currency_cell(tbl.cell(r, col_other_exp), other_exp)
        if col_total_exp is not None:
            _set_currency_cell(tbl.cell(r, col_total_exp), total_exp)

        if col_monthly is not None:
            _set_currency_cell(tbl.cell(r, col_monthly), monthly)
        if col_cumulative is not None:
            _set_currency_cell(tbl.cell(r, col_cumulative), cumulative_running)

        # Accumulate for totals row
        total_rent += rent
        total_other_rev += other_rev
        total_mortgage += mortgage
        total_hoa += hoa
        total_mgt_fee += mgt_fee
        total_repairs += repairs
        total_other_exp += other_exp
        total_monthly += monthly

    # Total row is last row: sum of everything (excluding timeframe = 'N/A' already handled in query)
    total_rev_all = total_rent + total_other_rev
    total_exp_all = total_mortgage + total_hoa + total_mgt_fee + total_repairs + total_other_exp
    total_monthly_all = total_rev_all + total_exp_all

    if col_rent is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_rent), total_rent)
    if col_other_rev is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_other_rev), total_other_rev)
    if col_total_rev is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_total_rev), total_rev_all)

    if col_mortgage is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_mortgage), total_mortgage)
    if col_hoa is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_hoa), total_hoa)
    if col_mgt_fee is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_mgt_fee), total_mgt_fee)
    if col_repairs is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_repairs), total_repairs)
    if col_other_exp is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_other_exp), total_other_exp)
    if col_total_exp is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_total_exp), total_exp_all)

    if col_monthly is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_monthly), total_monthly_all)
    if col_cumulative is not None:
        _set_currency_cell(tbl.cell(total_row_idx, col_cumulative), total_monthly_all)

    print("monthly_perf_table updated.")

OBJECT_UPDATERS: Dict[str, ObjectUpdater] = {
    "cover_title": _update_cover_title,
    "cover_subtitle": _update_cover_subtitle,
    "summary_table": _update_summary_table,
    "monthly_perf_table": _update_monthly_perf_table,
}
