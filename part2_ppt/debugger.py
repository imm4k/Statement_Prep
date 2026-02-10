from __future__ import annotations

import sys
from pathlib import Path
import zipfile

from pptx import Presentation


PPTX_PATH = Path(
    r"H:\.shortcut-targets-by-id\1Tf1JC85Wg2bbW79jfilWWqexA8jYAOxs"
    r"\CARTER Property Management\23. Operations and Administrative"
    r"\Property Management\Investor Updates\0. Templates"
    r"\TEMPLATE_Monthly_Nicholas Granato.pptx"
)

TARGET_SHAPE_NAME = "summary_top_text"
TOKENS = ["[CoC Return]", "[Cumulative Return]"]


def _zip_precheck(p: Path) -> None:
    with zipfile.ZipFile(p, "r") as z:
        bad = z.testzip()
        if bad is not None:
            raise RuntimeError(f"Zip integrity check failed. First bad entry: {bad}")


def _debug_shape_text_runs(pptx_path: Path) -> None:
    print(f"PPTX: {pptx_path}")

    if not pptx_path.exists():
        print("ERROR: File not found.")
        return

    try:
        _zip_precheck(pptx_path)
        print("Zip precheck: OK")
    except Exception as e:
        print(f"Zip precheck: FAILED ({type(e).__name__}: {e})")
        return

    try:
        prs = Presentation(str(pptx_path))
        print("python-pptx load: OK")
    except Exception as e:
        print(f"python-pptx load: FAILED ({type(e).__name__}: {e})")
        return

    found = False

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "name", None) != TARGET_SHAPE_NAME:
                continue

            found = True
            print(f"\nFound shape '{TARGET_SHAPE_NAME}' on slide {slide_idx}")

            if not shape.has_text_frame:
                print("Shape has no text frame.")
                return

            tf = shape.text_frame
            print(f"Paragraphs: {len(tf.paragraphs)}")

            # Print runs and detect tokens
            for p_idx, p in enumerate(tf.paragraphs, start=1):
                print(f"\nParagraph {p_idx}")
                print(f"Paragraph text: '{p.text}'")
                if not p.runs:
                    print("  (no runs)")
                    continue

                for r_idx, run in enumerate(p.runs, start=1):
                    raw = run.text
                    print(f"  Run {r_idx}: '{raw}'")
                    for token in TOKENS:
                        if token in raw:
                            print(f"    TOKEN FOUND IN SINGLE RUN: {token}")

            # Detect tokens across runs (handles split tokens)
            full_text = "".join(run.text for p in tf.paragraphs for run in p.runs)
            print("\nReconstructed full text across runs:")
            print(f"'{full_text}'")

            for token in TOKENS:
                if token in full_text:
                    print(f"TOKEN FOUND ACROSS RUNS: {token}")
                else:
                    print(f"TOKEN NOT FOUND: {token}")

            # Extra: show where token chars appear if split oddly
            for token in TOKENS:
                if token not in full_text:
                    continue
                idx = full_text.find(token)
                print(f"{token} position in full_text: index {idx}")

            return

    if not found:
        print(f"Shape '{TARGET_SHAPE_NAME}' not found.")
        # Helpful inventory (names only) so you can confirm Selection Pane naming
        print("\nAvailable shape names:")
        names = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                nm = getattr(shape, "name", "")
                if nm:
                    names.append((slide_idx, nm))
        for slide_idx, nm in names:
            if "summary" in nm.lower() or "top" in nm.lower() or "text" in nm.lower():
                print(f"  slide {slide_idx}: {nm}")


def main() -> None:
    # Optional override: run with an explicit path if you want
    pptx_path = PPTX_PATH
    if len(sys.argv) == 2:
        pptx_path = Path(sys.argv[1])

    _debug_shape_text_runs(pptx_path)


if __name__ == "__main__":
    main()
